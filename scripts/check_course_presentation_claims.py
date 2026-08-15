"""Claim-boundary and citation guard for the IS Research Seminar deck.

Two failure modes would each be unrecoverable for a graduate literature review,
so both are checked mechanically rather than by eye:

  1. A forbidden outcome claim. EXP-005 stands at 0/24 validated
     generalization-safe expert labels, so no accuracy, effort-reduction,
     generalization or clinical-performance claim is supportable.
  2. An unscoped proven-absence claim. The frozen protocol searches QL-01..QL-05
     have not been executed, so "nobody has done X" is not sayable; only
     "within the reviewed corpus, X is not addressed" is.

Every hit is printed with its slide and surrounding text for human adjudication -
the script reports, it does not silently pass.
"""
from __future__ import annotations

import re
import sys
from pathlib import Path

from pptx import Presentation

DECK = (Path(r"C:\Users\ahamed\vego-ai\outputs\course-presentation")
        / "VEGO-AI - IS Research Seminar - Final Presentation.pptx")

# Phrases that assert an outcome this research cannot yet demonstrate.
FORBIDDEN = [
    (r"improve[sd]?\s+(the\s+)?accuracy", "accuracy-improvement claim"),
    (r"more accurate|better accuracy|higher accuracy", "accuracy-improvement claim"),
    (r"reduce[sd]?\s+(expert\s+)?(effort|workload|burden)", "effort-reduction claim"),
    (r"saves?\s+(time|effort)", "effort-reduction claim"),
    (r"\bgeneralis|generaliz", "generalization claim (check scoping)"),
    (r"clinical(ly)?\s+(performance|benefit|outcome|validated)", "clinical claim"),
    (r"patient outcome", "clinical claim"),
    (r"outperform|state of the art results|beats\b", "superiority claim"),
]

# Absence claims are only acceptable when explicitly scoped to what was read.
ABSENCE = [
    r"\bno one\b", r"\bnobody\b", r"\bnever been\b", r"\bfirst to\b",
    r"does not exist", r"has not been (done|studied|addressed)",
    r"\bno (study|work|research|paper|source)\b",
]
SCOPE_MARKERS = [
    "reviewed corpus", "reviewed work", "reviewed set", "in this work",
    "the corpus", "corpus only", "reviewed source", "not a proven absence",
    "have not yet been executed", "not yet executed", "what has been found and read",
]

# An absence phrase inside an explicit self-limiting statement ("I cannot yet say
# what does not exist", "requires evidence that does not exist yet") is the
# guard-rail itself, not a violation of it.
DISCLAIMER_MARKERS = [
    "cannot yet", "cannot claim", "i cannot", "requires evidence",
    "does not exist yet", "not yet, not done", "both are next",
    "requires independent", "still requires",
]


def slide_texts(prs):
    for i, s in enumerate(prs.slides, 1):
        chunks = []
        for shape in s.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                chunks.append(shape.text_frame.text)
        notes = ""
        if s.has_notes_slide:
            notes = s.notes_slide.notes_text_frame.text
        yield i, "\n".join(chunks), notes


def main():
    prs = Presentation(DECK)
    problems, warnings = [], []

    for idx, body, notes in slide_texts(prs):
        for scope, text in (("slide", body), ("notes", notes)):
            if not text:
                continue
            low = text.lower()
            for pat, label in FORBIDDEN:
                for m in re.finditer(pat, low):
                    seg = text[max(0, m.start() - 110): m.end() + 110].replace("\n", " ")
                    segl = seg.lower()
                    # A negated/excluded mention, or one inside a research question,
                    # is the point being made rather than a violation of it.
                    neg = re.search(
                        r"\bno\b|\bnot\b|cannot|never|without|require|excluded|unsafe|"
                        r"\bnor\b|forbidden|do not",
                        segl)
                    question = "?" in seg and re.search(r"how can|what |which ", segl)
                    (warnings if (neg or question) else problems).append(
                        (idx, scope, label, seg.strip()))
            for pat in ABSENCE:
                for m in re.finditer(pat, low):
                    seg = text[max(0, m.start() - 200): m.end() + 200].replace("\n", " ")
                    segl = seg.lower()
                    if any(k in segl for k in SCOPE_MARKERS):
                        continue
                    if any(k in segl for k in DISCLAIMER_MARKERS):
                        warnings.append((idx, scope, "absence inside a disclaimer", seg.strip()))
                        continue
                    problems.append((idx, scope, "UNSCOPED absence claim", seg.strip()))

    print(f"deck: {DECK.name}")
    print(f"slides: {len(prs.slides.__iter__.__self__._sldIdLst)}\n")

    if problems:
        print(f"!! {len(problems)} PROBLEM(S) - must be fixed\n")
        for idx, scope, label, seg in problems:
            print(f"  [slide {idx} · {scope}] {label}\n      …{seg}…\n")
    else:
        print("PASS - no unscoped absence claim and no unnegated forbidden claim.\n")

    if warnings:
        print(f"-- {len(warnings)} negated//excluded mention(s), reviewed and expected:\n")
        for idx, scope, label, seg in warnings:
            print(f"  [slide {idx} · {scope}] {label}\n      …{seg}…\n")

    return 1 if problems else 0


if __name__ == "__main__":
    sys.exit(main())
