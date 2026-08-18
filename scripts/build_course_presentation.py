"""Build the IS Research Seminar (214.4001) final presentation.

Course: "The Research Approach in Information Systems", Prof. Penina Soffer.
Assignment (CL7 slide 11): a 10-12 minute class presentation covering motivation,
the research question, the derived set of questions, the literature search
strategy, and initial findings (research streams, classification/analysis
framework, identified gaps).

Content sources - every one of them a real artifact in this repository:
  * literature/verified-research-corpus-2026-08-12.json  (144 verified sources)
  * reports/generated/exp006|exp007|exp008/summary.json  (measured mechanism runs)
  * docs/research/governance/vego-ai-foundation-paper-record.md (published figures)
  * docs/research/phd-proposal/three-study-contract.md   (RQ / SQ wording)
  * artifacts/meetings/2026-08-12-iris-arnon/            (supervisor directives)
  * outputs/course-presentation/findings.json            (analysis workflow output)

Evidence boundary held throughout: EXP-005 stands at 0/24 validated
generalization-safe expert labels, and the frozen searches QL-01..QL-05 have not
been executed. Therefore the deck makes no accuracy, effort-reduction,
generalization or clinical claim, and never asserts proven absence in the
literature - only what the reviewed corpus does and does not cover.
"""
from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

REPO = Path(r"C:\Users\ahamed\vego-ai")
OUT_DIR = REPO / "outputs" / "course-presentation"
FIG = OUT_DIR / "figures"
OUT_DIR.mkdir(parents=True, exist_ok=True)
DECK = OUT_DIR / "VEGO-AI - IS Research Seminar - Final Presentation.pptx"

# ---------------------------------------------------------------- design system
NAVY = RGBColor(0x0F, 0x27, 0x40)
NAVY_SOFT = RGBColor(0x1B, 0x3A, 0x5A)
ACCENT = RGBColor(0x2A, 0x78, 0xD6)
ACCENT_L = RGBColor(0xCD, 0xE2, 0xFB)
ORANGE = RGBColor(0xEB, 0x68, 0x34)
CRIT = RGBColor(0xD0, 0x3B, 0x3B)
GOOD = RGBColor(0x0C, 0xA3, 0x0C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xF2, 0xF5, 0xF9)
SOFT2 = RGBColor(0xE8, 0xEE, 0xF6)
INK = RGBColor(0x0B, 0x0B, 0x0B)
MUTED = RGBColor(0x5A, 0x64, 0x72)
LINE = RGBColor(0xD5, 0xDD, 0xE7)

H_FONT = "Cambria"     # safe-list serif for headings
B_FONT = "Calibri"     # safe-list sans for body

W, H = Inches(13.333), Inches(7.5)
M = Inches(0.62)                     # page margin
CONTENT_W = W - 2 * M


def new_deck() -> Presentation:
    p = Presentation()
    p.slide_width, p.slide_height = W, H
    return p


def slide(prs, dark=False):
    s = prs.slides.add_slide(prs.slide_layouts[6])   # blank
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY if dark else WHITE
    return s


def tb(slide, x, y, w, h, *, anchor=MSO_ANCHOR.TOP, wrap=True):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return box, tf


def para(tf, text, *, size=16, bold=False, color=INK, font=B_FONT, align=PP_ALIGN.LEFT,
         space_after=0, space_before=0, italic=False, line=None, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    if line:
        p.line_spacing = line
    r = p.add_run()
    r.text = text
    f = r.font
    f.size, f.bold, f.italic, f.name = Pt(size), bold, italic, font
    f.color.rgb = color
    return p


def rich(tf, chunks, *, size=16, align=PP_ALIGN.LEFT, space_after=0, space_before=0,
         line=None, first=False):
    """A paragraph made of differently-styled runs: [(text, {bold,color,...}), ...]"""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    if line:
        p.line_spacing = line
    for text, st in chunks:
        r = p.add_run()
        r.text = text
        f = r.font
        f.size = Pt(st.get("size", size))
        f.bold = st.get("bold", False)
        f.italic = st.get("italic", False)
        f.name = st.get("font", B_FONT)
        f.color.rgb = st.get("color", INK)
    return p


def card(slide, x, y, w, h, *, fill=SOFT, line_col=None, radius=0.04, shadow=False):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = radius
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line_col:
        sh.line.color.rgb = line_col
        sh.line.width = Pt(1.0)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = shadow
    sh.text_frame.word_wrap = True
    return sh


def badge(slide, x, y, d, text, *, fill=ACCENT, color=WHITE, size=15):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size, r.font.bold, r.font.name = Pt(size), True, B_FONT
    r.font.color.rgb = color
    return sh


# Rough advance width of Cambria bold, as a fraction of point size. Used only to
# predict how many lines a heading will wrap to, so the subtitle can be placed
# below it instead of on top of it.
_TITLE_CHAR_W = 0.47
_SUB_CHAR_W = 0.47


def _wrapped_lines(text, pt_size, width_in, char_w):
    """Estimate rendered line count for `text` at `pt_size` inside `width_in`."""
    per_line = max(12, int(width_in / (pt_size * char_w / 72.0)))
    lines, cur = 1, 0
    for word in text.split():
        add = len(word) + (1 if cur else 0)
        if cur + add > per_line:
            lines += 1
            cur = len(word)
        else:
            cur += add
    return lines


def heading(slide, title, *, eyebrow=None, dark=False, sub=None, size=31):
    y = Inches(0.46)
    if eyebrow:
        _, tf = tb(slide, M, y, CONTENT_W, Inches(0.26))
        para(tf, eyebrow.upper(), size=11.5, bold=True,
             color=ACCENT if not dark else ACCENT_L, first=True)
        y += Inches(0.32)

    w_in = CONTENT_W / 914400
    n = _wrapped_lines(title, size, w_in, _TITLE_CHAR_W)
    line_h = Inches(size / 72.0 * 1.22)
    _, tf = tb(slide, M, y, CONTENT_W, line_h * n + Inches(0.08))
    para(tf, title, size=size, bold=True, color=WHITE if dark else NAVY, font=H_FONT,
         first=True, line=1.1)
    y += line_h * n + Inches(0.16)

    if sub:
        ns = _wrapped_lines(sub, 14.5, w_in, _SUB_CHAR_W)
        sub_h = Inches(14.5 / 72.0 * 1.30) * ns
        _, tf = tb(slide, M, y, CONTENT_W, sub_h + Inches(0.06))
        para(tf, sub, size=14.5, color=ACCENT_L if dark else MUTED, first=True, line=1.18)
        y += sub_h + Inches(0.22)
    return y


def footnote(slide, text, *, dark=False, y=None):
    y = y or (H - Inches(0.62))
    _, tf = tb(slide, M, y, CONTENT_W, Inches(0.42))
    para(tf, text, size=10.5, color=RGBColor(0x8F, 0x9B, 0xA8) if dark else MUTED,
         first=True, line=1.15)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()


def picture(slide, path, x, y, w=None, h=None):
    return slide.shapes.add_picture(str(path), x, y, width=w, height=h)


def fit_picture(slide, path, box_x, box_y, box_w, box_h):
    """Insert a picture scaled to fit inside a box, centred."""
    from PIL import Image
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(box_w / iw, box_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    x = box_x + (box_w - w) // 2
    y = box_y + (box_h - h) // 2
    return slide.shapes.add_picture(str(path), Emu(int(x)), Emu(int(y)),
                                    width=Emu(w), height=Emu(h))


def arrow(slide, x, y, w, h, *, color=ACCENT):
    sh = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


# ---------------------------------------------------------------- content
F = json.loads((OUT_DIR / "findings.json").read_text(encoding="utf-8"))

U_RQ = ("How can human judgment be captured, governed, and used to support "
        "agentic-AI-driven variability exploration in guideline operationalization "
        "scenarios, enabling reliable human–AI co-reasoning?")

SQS = [
    ("SQ1", "When to ask",
     "When and how should an agentic assessment system request human judgment, so that "
     "important uncertainties are addressed without unnecessary expert burden?",
     "Study 1 · selective intervention"),
    ("SQ2", "What to capture and govern",
     "How should expert judgment — including the system's own reasoning — be represented, "
     "validated, reconciled and stored so it can be reused transparently, without unsafe "
     "generalization or loss of human authority?",
     "Study 2 · judgment lifecycle"),
    ("SQ3", "When to reuse and transfer",
     "How can expert judgment be reused and transferred across guideline-operationalization "
     "contexts — what generalizes, and what must adapt?",
     "Study 3 · transfer and evaluation"),
]


def build():
    prs = new_deck()

    # ======================================================= 1 TITLE
    s = slide(prs, dark=True)
    _, tf = tb(s, M, Inches(0.98), CONTENT_W, Inches(0.3))
    para(tf, "IS RESEARCH SEMINAR  ·  214.4001  ·  FINAL PRESENTATION",
         size=12, bold=True, color=ACCENT_L, first=True)

    _, tf = tb(s, M, Inches(1.46), Inches(11.4), Inches(1.5))
    para(tf, "Not all differences matter.", size=41, bold=True, color=WHITE,
         font=H_FONT, first=True, line=1.06)
    para(tf, "So who decides which ones do?", size=41, bold=True, color=ACCENT_L,
         font=H_FONT, line=1.06)

    _, tf = tb(s, M, Inches(3.06), Inches(11.2), Inches(0.5))
    para(tf, "Governed human judgment in agentic-AI variability exploration — "
             "a structured exploratory literature review",
         size=16, color=RGBColor(0xC8, 0xD6, 0xE6), first=True, line=1.2)

    card(s, M, Inches(3.78), Inches(11.4), Inches(1.40), fill=NAVY_SOFT)
    _, tf = tb(s, M + Inches(0.34), Inches(3.96), Inches(10.75), Inches(1.05))
    para(tf, "RESEARCH QUESTION", size=10.5, bold=True, color=ACCENT_L, first=True,
         space_after=5)
    para(tf, U_RQ, size=14, color=WHITE, italic=True, line=1.22)

    _, tf = tb(s, M, Inches(5.44), Inches(11.4), Inches(1.2))
    rich(tf, [("Ali Hamed", {"bold": True, "size": 15, "color": WHITE})], first=True,
         space_after=4)
    para(tf, "M.Sc. Information Systems · University of Haifa", size=12.5,
         color=RGBColor(0xA9, 0xBC, 0xD0))
    para(tf, "Supervisors: Prof. Iris Reinhartz-Berger · Prof. Arnon Sturm",
         size=12.5, color=RGBColor(0xA9, 0xBC, 0xD0), space_before=2)
    para(tf, "Doctoral working draft · research-question wording provisional, "
             "not supervisor-approved · no empirical or medical claim.",
         size=10.5, color=RGBColor(0x7E, 0x94, 0xAC), italic=True, space_before=6)
    notes(s, """
Open with the boundary, not a boast. This is a controlled literature review and a research
agenda - not an empirical result.

Then the tension: Iris's own MODELS paper is called "Not All Differences Matter" - the class
saw her present it in week 5. My question starts one step later. Once an agent has found a
difference, WHO decides whether it is a legitimate alternative or an error, and how does that
decision stop being thrown away?  (~35 sec)
""")

    # ======================================================= 2 MOTIVATION
    s = slide(prs)
    y = heading(s, "An AI can find every difference — not which ones matter",
                eyebrow="Motivation · the problem",
                sub="The same structural deviation can be three different things, and only "
                    "situated judgment separates them.")
    cw, gap = Inches(3.72), Inches(0.28)
    items = [
        ("A", "A legitimate alternative", GOOD,
         "A different representation can be valid under the same language semantics and "
         "domain meaning. Penalising it teaches the wrong lesson."),
        ("B", "A genuine error", CRIT,
         "A misconception, omission, syntax violation or domain mistake that should be "
         "corrected."),
        ("C", "A defective guideline", ORANGE,
         "The reference itself may be incomplete, underspecified or mismatched to the "
         "case — the deviation is evidence against the norm."),
    ]
    for i, (letter, title, col, body) in enumerate(items):
        x = M + i * (cw + gap)
        card(s, x, y, cw, Inches(2.62), fill=SOFT, line_col=LINE)
        badge(s, x + Inches(0.30), y + Inches(0.30), Inches(0.50), letter, fill=col)
        _, tf = tb(s, x + Inches(0.30), y + Inches(1.00), cw - Inches(0.60), Inches(1.5))
        para(tf, title, size=16.5, bold=True, color=NAVY, font=H_FONT, first=True,
             space_after=8)
        para(tf, body, size=13, color=MUTED, line=1.26)

    yb = y + Inches(2.94)
    card(s, M, yb, CONTENT_W, Inches(1.00), fill=NAVY)
    _, tf = tb(s, M + Inches(0.34), yb + Inches(0.18), CONTENT_W - Inches(0.68), Inches(0.70))
    para(tf, "“Human-in-the-loop” is not one problem. It is four design problems:",
         size=16.5, bold=True, color=WHITE, font=H_FONT, first=True, space_after=5)
    para(tf, "when to ask  ·  what to capture  ·  how to govern  ·  when to reuse",
         size=14, color=ACCENT_L)
    notes(s, """
Ground the problem before any technology. The same difference can be a valid alternative, a
real error, or evidence that the guideline itself is wrong.

Then the reframing that drives the whole review: human oversight is not a slogan or an
approval button. It decomposes into four testable design problems. Those four become my
sub-questions.  (~50 sec)
""")

    # ======================================================= 3 PUBLISHED EVIDENCE
    s = slide(prs)
    y = heading(s, "The framework names its own need for a human",
                eyebrow="Motivation · evidence from the problem world",
                sub="VEGO-AI assesses domain models with four coordinated LLM agents. "
                    "It works — unevenly, and it says so.")
    fit_picture(s, FIG / "04-published-profile-bare.png", M, y - Inches(0.02),
                Inches(8.30), Inches(4.20))
    x2 = M + Inches(8.52)
    cw2 = CONTENT_W - Inches(8.52)
    card(s, x2, y - Inches(0.02), cw2, Inches(4.20), fill=SOFT, line_col=LINE)
    _, tf = tb(s, x2 + Inches(0.26), y + Inches(0.20), cw2 - Inches(0.52), Inches(3.8))
    para(tf, "ρ = 0.22", size=32, bold=True, color=CRIT, font=H_FONT, first=True)
    para(tf, "Reported correlation between Model-Inspector compliance scores and the "
             "human grader.",
         size=12, color=MUTED, line=1.22, space_before=3, space_after=12)
    para(tf, "The paper's own future work:", size=12, bold=True, color=NAVY, space_after=5)
    para(tf, "“incorporating human-in-the-loop oversight at key pipeline stages "
             "(guideline validation, compliance review, and variability classification) "
             "would strengthen reliability”",
         size=12, color=NAVY, italic=True, line=1.24, space_after=11)
    para(tf, "The need for human judgment is not my assumption. It is the framework's own "
             "stated next step — and it names the three stages.",
         size=12, bold=True, color=INK, line=1.24)
    footnote(s, "Figures as reported by the paper for its own education-domain evaluation "
                "(two domains, one institution, one LLM, two modelling languages) — "
                "reported findings, not independently reproduced here.")
    notes(s, """
This is the strongest motivation slide because almost none of it is mine.

The framework's own authors report that uncovered-fragment auditing sits at 0.55 in both
use-case settings, and that agreement with the human grader is weak. Then their future-work
section names exactly the three pipeline stages where oversight is needed - which is my SQ1,
in someone else's peer-reviewed paper.

Be precise: these are THEIR numbers, in THEIR setting, reported not reproduced. I claim
nothing about my own performance.  (~55 sec)
""")

    # ======================================================= 4 MEASURED EVIDENCE
    s = slide(prs)
    y = heading(s, "The same gap is measurable inside the pipeline",
                eyebrow="Motivation · evidence from my own instrumentation",
                sub="Replaying the assessment pipeline as an event stream shows how little "
                    "of its reasoning was ever reviewable — and that asking more often "
                    "does not fix it.")
    fit_picture(s, FIG / "02-observability-gap-bare.png", M, y, Inches(7.80), Inches(3.60))
    x2 = M + Inches(8.05)
    cw2 = CONTENT_W - Inches(8.05)
    for i, (big, lbl, col) in enumerate([
        ("44×", "more lifecycle events occurred than the legacy review queue could show",
         ACCENT),
        ("33", "reference guidelines per setting revised with no human ever seeing them",
         ORANGE),
        ("0.80+", "share of events a human must see to reach 80% high-severity coverage",
         CRIT),
    ]):
        cy = y + i * Inches(1.26)
        card(s, x2, cy, cw2, Inches(1.10), fill=SOFT, line_col=LINE)
        _, tf = tb(s, x2 + Inches(0.24), cy + Inches(0.14), cw2 - Inches(0.48), Inches(0.86))
        para(tf, big, size=25, bold=True, color=col, font=H_FONT, first=True, space_after=2)
        para(tf, lbl, size=11, color=MUTED, line=1.16)
    footnote(s, "EXP-006/007/008 · mechanism and observability evidence only. These are "
                "heterogeneous lifecycle observations and routing trade-offs, not quality "
                "outcomes — they support no claim about assessment accuracy.")
    notes(s, """
The published evaluation looks at outputs. I instrumented the process.

Four settings produced 481 lifecycle events; the review queue a human actually had could
surface 11. Thirty-three reference guidelines per setting were revised without any human
review - and those guidelines are what every assessment is measured against.

Then the third number, which is the real research problem: to catch 80% of high-severity
events you must route roughly 80% of them to a person. That defeats the automation. So "add
a human in the loop" is not yet an answer - WHICH cases, WHEN, and at what dose is open.

Careful: this says the process was invisible and the trade-off is real. It does NOT say the
assessments were wrong. I have no expert labels yet.  (~55 sec)
""")

    # ======================================================= 5 BASELINE
    s = slide(prs)
    y = heading(s, "Where judgment would enter the pipeline",
                eyebrow="Research context · the VEGO-AI baseline",
                sub="Four coordinated agents. The review asks what a governed human-judgment "
                    "layer around them would have to do.")
    agents = [
        ("Language Advisor", "language template · metamodel semantics",
         "language-rule ambiguity"),
        ("Domain Advisor", "reference guidelines · valid alternatives",
         "guideline mapping, scope"),
        ("Model Inspector", "case assessment · uncovered fragments",
         "disputed compliance"),
        ("Variability Explorer", "recurring patterns · population view",
         "pattern interpretation"),
    ]
    # Fill the band between heading and footnote instead of fixed heights, so a
    # wrapped heading never leaves a dead strip at the bottom.
    avail5 = (H - Inches(0.72)) - y
    h_agent = avail5 * 0.34
    h_need = avail5 * 0.18
    h_bot = avail5 - h_agent - h_need - Inches(0.44)
    aw = (CONTENT_W - Inches(0.24) * 3) / 4
    for i, (name, does, needs) in enumerate(agents):
        x = M + i * (aw + Inches(0.24))
        card(s, x, y, aw, h_agent, fill=NAVY)
        _, tf = tb(s, x + Inches(0.20), y + Inches(0.18), aw - Inches(0.40),
                   h_agent - Inches(0.36), anchor=MSO_ANCHOR.MIDDLE)
        para(tf, name, size=15, bold=True, color=WHITE, font=H_FONT, first=True,
             space_after=8, line=1.08)
        para(tf, does, size=11, color=ACCENT_L, line=1.24)
        cy2 = y + h_agent + Inches(0.16)
        card(s, x, cy2, aw, h_need, fill=RGBColor(0xFD, 0xF0, 0xF0),
             line_col=RGBColor(0xF0, 0xC8, 0xC8))
        _, tf = tb(s, x + Inches(0.18), cy2 + Inches(0.06), aw - Inches(0.36),
                   h_need - Inches(0.12), anchor=MSO_ANCHOR.MIDDLE)
        para(tf, needs, size=11, bold=True, color=CRIT, first=True, line=1.14)

    yb = y + h_agent + h_need + Inches(0.44)
    lw = (CONTENT_W - Inches(0.24)) / 2
    card(s, M, yb, lw, h_bot, fill=SOFT, line_col=LINE)
    _, tf = tb(s, M + Inches(0.26), yb + Inches(0.18), lw - Inches(0.52), h_bot - Inches(0.36))
    para(tf, "Reported as working", size=13, bold=True, color=GOOD, font=H_FONT,
         first=True, space_after=6)
    para(tf, "Modular separation of language and domain concerns; reasonably stable "
             "templates and guidelines in the reported settings.",
         size=11.5, color=MUTED, line=1.2)
    card(s, M + lw + Inches(0.24), yb, lw, h_bot, fill=RGBColor(0xFD, 0xF0, 0xF0),
         line_col=RGBColor(0xF0, 0xC8, 0xC8))
    _, tf = tb(s, M + lw + Inches(0.50), yb + Inches(0.18), lw - Inches(0.52),
               h_bot - Inches(0.36))
    para(tf, "Not established", size=13, bold=True, color=CRIT, font=H_FONT, first=True,
         space_after=6)
    para(tf, "No general correctness across institutions, LLMs or domains; no validated "
             "reusable human-judgment mechanism; no healthcare readiness.",
         size=11.5, color=MUTED, line=1.2)
    notes(s, """
Brief research context, not a product demo. Four agents, each producing something a human
might need to rule on - the red boxes are the insertion points.

Then the honest split at the bottom. What the manuscript reports as working, and what is
simply not established. I am building on their foundation, not claiming their results as
mine - and I am not an author of that paper.  (~50 sec)
""")

    # ======================================================= 6 RESEARCH QUESTION
    s = slide(prs, dark=True)
    y = heading(s, "The research question", eyebrow="What I am asking", dark=True)
    card(s, M, y + Inches(0.10), CONTENT_W, Inches(1.46), fill=NAVY_SOFT)
    _, tf = tb(s, M + Inches(0.4), y + Inches(0.34), CONTENT_W - Inches(0.8), Inches(1.05))
    para(tf, U_RQ, size=19.5, color=WHITE, font=H_FONT, italic=True, first=True, line=1.24)

    y2 = y + Inches(1.86)
    cw3, gap3 = Inches(3.86), Inches(0.24)
    commitments = [
        ("It names the task, not the solution",
         "The object of study is variability exploration in guideline operationalization — "
         "not a particular agent architecture."),
        ("It is narrower than “human–AI collaboration”",
         "Intervention and reuse must attach to a real decision unit: a guideline mapping, "
         "a fragment classification, a pattern."),
        ("Reuse is not assumed in the headline",
         "Whether captured judgment may be reused is a governed, testable concern for the "
         "sub-questions — not a premise."),
    ]
    for i, (t, b) in enumerate(commitments):
        x = M + i * (cw3 + gap3)
        card(s, x, y2, cw3, Inches(1.86), fill=NAVY_SOFT)
        _, tf = tb(s, x + Inches(0.26), y2 + Inches(0.22), cw3 - Inches(0.52), Inches(1.44))
        para(tf, t, size=14, bold=True, color=ACCENT_L, font=H_FONT, first=True,
             space_after=7, line=1.14)
        para(tf, b, size=12, color=RGBColor(0xBF, 0xCE, 0xDE), line=1.24)
    footnote(s, "Drafting baseline, not supervisor-approved wording. D-RQ-01/D-RQ-02 remain "
                "open, along with “exploration” versus “identification and "
                "classification” and “human” versus “expert” judgment.",
             dark=True)
    notes(s, """
Read the question once, slowly. Then the three decisions behind its wording, because this
course cares about how a question is built.

It names the task not the solution. It is deliberately narrower than "human-AI collaboration"
- it has to attach to a real decision unit. And it keeps reuse OUT of the headline, so reuse
stays something I test rather than assume.

Be honest that the wording is still provisional; two specific words are unresolved.  (~45 sec)
""")

    # ======================================================= 7 SUB-QUESTIONS
    s = slide(prs)
    y = heading(s, "Three derived questions — which become three studies",
                eyebrow="Derived set of questions",
                sub="Each sub-question owns one stage: when to ask, what to keep and govern, "
                    "and where a ruling may travel.")
    ch = Inches(1.32)
    for i, (key, title, body, study) in enumerate(SQS):
        cy = y + i * (ch + Inches(0.18))
        card(s, M, cy, CONTENT_W, ch, fill=SOFT, line_col=LINE)
        badge(s, M + Inches(0.30), cy + Inches(0.28), Inches(0.60), key, size=14)
        _, tf = tb(s, M + Inches(1.12), cy + Inches(0.18), Inches(7.5), Inches(1.0))
        para(tf, title, size=16, bold=True, color=NAVY, font=H_FONT, first=True, space_after=5)
        para(tf, body, size=12, color=MUTED, line=1.22)
        card(s, M + Inches(8.86), cy + Inches(0.24), Inches(3.08), Inches(0.84),
             fill=ACCENT_L)
        _, tf = tb(s, M + Inches(9.00), cy + Inches(0.38), Inches(2.80), Inches(0.58),
                   anchor=MSO_ANCHOR.MIDDLE)
        para(tf, study, size=12.5, bold=True, color=NAVY, first=True, line=1.16)
    footnote(s, "Separating the studies prevents one small label set from being treated as "
                "evidence for intervention policy, representation validity and transfer at "
                "once. The umbrella question is answered only after they integrate.")
    notes(s, """
Three sub-questions, and each has to become a study with its own artifact, comparator and
falsification criterion - that is the structure my supervisors asked for, and it maps onto
the design-science framing from this course.

The footnote matters: separating them is a methodological safeguard. Otherwise the same
twenty-odd expert labels get reused as evidence for everything, which would be indefensible.

Say it plainly: the questions are domain-neutral. Software engineering and medicine are where
I would test them, not what they are about.  (~55 sec)
""")

    # ======================================================= 8 METHOD
    s = slide(prs)
    y = heading(s, "How the review was done",
                eyebrow="Method",
                sub="A structured exploratory review with systematic controls — organised "
                    "by concepts, not by paper summaries.")
    left_w = Inches(7.05)
    card(s, M, y, left_w, Inches(3.86), fill=SOFT, line_col=LINE)
    _, tf = tb(s, M + Inches(0.30), y + Inches(0.22), left_w - Inches(0.60), Inches(3.45))
    para(tf, "Relevant research areas", size=15, bold=True, color=NAVY, font=H_FONT,
         first=True, space_after=8)
    for t, d in [
        ("Human involvement in agentic AI — the core",
         "mixed initiative, levels of agency, oversight, control"),
        ("Selective intervention",
         "learning to defer, abstention, uncertainty, value of information, burden"),
        ("Judgment capture and governance",
         "reasoning capture, provenance, disagreement, authority, expiry, revocation"),
        ("Memory, reuse and transfer",
         "scope-aware retrieval, domain and guideline adaptation, leakage-safe evaluation"),
    ]:
        rich(tf, [(f"{t}  ", {"bold": True, "size": 12.5, "color": INK}),
                  (d, {"size": 12.5, "color": MUTED})], space_after=7, line=1.2)
    para(tf, "Deliberately outside the review body: enabling technologies (model "
             "architectures, local-inference tooling). They are engineering means — the "
             "review must justify the problem, not describe the solution.",
         size=11.5, color=MUTED, italic=True, space_before=5, line=1.2)

    rx = M + left_w + Inches(0.26)
    rw = CONTENT_W - left_w - Inches(0.26)
    card(s, rx, y, rw, Inches(1.82), fill=WHITE, line_col=LINE)
    _, tf = tb(s, rx + Inches(0.26), y + Inches(0.18), rw - Inches(0.52), Inches(1.5))
    para(tf, "Screening", size=14.5, bold=True, color=NAVY, font=H_FONT, first=True,
         space_after=7)
    rich(tf, [("Include  ", {"bold": True, "size": 12, "color": GOOD}),
              ("methods, frameworks, artifacts, empirical studies.",
               {"size": 12, "color": MUTED})], line=1.2, space_after=5)
    rich(tf, [("Exclude  ", {"bold": True, "size": 12, "color": CRIT}),
              ("duplicates, passing mentions, off-topic enabling tech, unverifiable sources.",
               {"size": 12, "color": MUTED})], line=1.2)

    card(s, rx, y + Inches(2.04), rw, Inches(1.82), fill=RGBColor(0xF3, 0xF8, 0xF3),
         line_col=RGBColor(0xCC, 0xE4, 0xCC))
    _, tf = tb(s, rx + Inches(0.26), y + Inches(2.22), rw - Inches(0.52), Inches(1.5))
    para(tf, "GenAI disclosure", size=14.5, bold=True, color=GOOD, font=H_FONT, first=True,
         space_after=7)
    para(tf, "Used for organisation, synthesis drafting, diagrams and consistency checks. "
             "Not used to approve inclusion, invent citations, or create expert labels.",
         size=12, color=MUTED, line=1.2)
    notes(s, """
The course asks for the search process to be explicit, so this is the honest version.

Four research areas. The first is the centre of gravity, and that was a deliberate
supervisory decision: organise the review around human involvement in agentic AI generally,
NOT around our specific application. The literature has to justify the problem, not describe
my solution.

Note what I exclude - enabling technology. Interesting, but it belongs in the methodology
chapter, not in the review that establishes the gap.

The GenAI box is a course requirement and I would rather show it than be asked.  (~55 sec)
""")

    # ======================================================= 9 SEARCH STRATEGY
    s = slide(prs)
    y = heading(s, "Search strategy — and what has actually been executed",
                eyebrow="Method · sources, queries and status",
                sub="The query families are locked. Presenting them is not the same as "
                    "claiming the formal searches have been run.")
    avail9 = (H - Inches(0.72)) - y
    fit_picture(s, FIG / "03-corpus-composition-bare.png", M, y, Inches(7.30), avail9)
    x2 = M + Inches(7.55)
    cw2 = CONTENT_W - Inches(7.55)
    h_q = avail9 * 0.56
    h_stat = avail9 - h_q - Inches(0.20)
    card(s, x2, y, cw2, h_q, fill=SOFT, line_col=LINE)
    _, tf = tb(s, x2 + Inches(0.24), y + Inches(0.20), cw2 - Inches(0.48), h_q - Inches(0.40))
    para(tf, "Locked query family — SQ1", size=13, bold=True, color=NAVY, font=H_FONT,
         first=True, space_after=6)
    para(tf, "(learning to defer OR algorithmic triage OR abstention OR selective "
             "prediction) AND (human OR expert) AND (cost OR burden OR timing OR "
             "uncertainty OR disagreement)",
         size=11, color=MUTED, line=1.24, space_after=8)
    para(tf, "Sources: Scopus · ACM DL · IEEE Xplore · SpringerLink · "
             "Web of Science.  Scholar and dblp for discovery, snowballing and verification.",
         size=11, color=MUTED, line=1.24)

    card(s, x2, y + h_q + Inches(0.20), cw2, h_stat, fill=RGBColor(0xFD, 0xF0, 0xF0),
         line_col=RGBColor(0xF0, 0xC8, 0xC8))
    _, tf = tb(s, x2 + Inches(0.24), y + h_q + Inches(0.36), cw2 - Inches(0.48),
               h_stat - Inches(0.36))
    para(tf, "Execution status", size=13, bold=True, color=CRIT, font=H_FONT, first=True,
         space_after=6)
    rich(tf, [("QL-01–QL-05:  ", {"size": 12, "color": MUTED}),
              ("0 of 5 executed", {"bold": True, "size": 12, "color": CRIT})],
         line=1.22, space_after=4)
    para(tf, "20 RQ-anchor mappings verified (5 per question); evidence maturity 14 full-text "
             "+ 6 record-level. This is what has been found and read — not evidence that "
             "nothing else exists.",
         size=11, color=MUTED, line=1.22)
    notes(s, """
This is the CL7 search-strategy requirement, answered without pretending.

The query families are locked and registered - I can show you the exact strings. The sources
are named. But the formal database searches have NOT been executed: zero of five.

So what IS the evidence base? Twenty anchor mappings, five per research question, each
bibliographically verified; fourteen at full text, six at record level.

That means I can say what I have read. I cannot say what does not exist. Every gap on the
next slides is stated inside that limit - and I will say so again when I get there.  (~50 sec)
""")

    # ======================================================= 10 STREAMS
    s = slide(prs)
    y = heading(s, F["streams_title"], eyebrow="Findings · research streams",
                sub=F["streams_sub"])
    rows = F["streams"]
    cols, gapx, gapy = 3, Inches(0.22), Inches(0.18)
    cw4 = (CONTENT_W - gapx * (cols - 1)) / cols
    avail = (H - Inches(0.78)) - y
    chh = (avail - gapy) / 2
    for i, r in enumerate(rows):
        cx = M + (i % cols) * (cw4 + gapx)
        cy = y + (i // cols) * (chh + gapy)
        card(s, cx, cy, cw4, chh, fill=SOFT, line_col=LINE)
        badge(s, cx + Inches(0.22), cy + Inches(0.18), Inches(0.38), r["key"], size=11.5)
        _, tf = tb(s, cx + Inches(0.68), cy + Inches(0.18), cw4 - Inches(0.92), Inches(0.40),
                   anchor=MSO_ANCHOR.MIDDLE)
        rich(tf, [(r["name"], {"bold": True, "size": 12.5, "color": NAVY, "font": H_FONT}),
                  ("   " + r["count"], {"bold": True, "size": 11, "color": ACCENT})],
             first=True, line=1.06)
        _, tf = tb(s, cx + Inches(0.22), cy + Inches(0.66), cw4 - Inches(0.44), Inches(1.5))
        para(tf, r["establishes"], size=10, color=MUTED, line=1.16, first=True,
             space_after=5)
        rich(tf, [("Leaves open  ", {"bold": True, "size": 10, "color": CRIT}),
                  (r["leaves_open"], {"size": 10, "color": MUTED})], line=1.16)
    footnote(s, F["streams_foot"])
    notes(s, F["streams_notes"])

    # ======================================================= 11 FRAMEWORK
    s = slide(prs)
    y = heading(s, F["framework_name"], eyebrow="Findings · classification framework",
                sub=F["framework_sub"])
    avail_h = (H - Inches(0.86)) - y
    fig_w = Inches(8.45)
    fit_picture(s, FIG / "05-maturity-grid-bare.png", M, y, fig_w, avail_h)
    px = M + fig_w + Inches(0.22)
    pw = CONTENT_W - fig_w - Inches(0.22)
    _, tf = tb(s, px, y, pw, Inches(0.28))
    para(tf, "WHAT THE COLUMNS REVEAL", size=11, bold=True, color=CRIT, first=True)
    top = y + Inches(0.36)
    gap_c = Inches(0.16)
    ch = ((H - Inches(0.86)) - top - gap_c * 2) / 3
    for i, (t, b) in enumerate(F["ceilings"]):
        cy = top + i * (ch + gap_c)
        card(s, px, cy, pw, ch, fill=SOFT, line_col=LINE)
        _, tf = tb(s, px + Inches(0.20), cy + Inches(0.14), pw - Inches(0.40), ch - Inches(0.24))
        para(tf, t, size=12.5, bold=True, color=NAVY, font=H_FONT, first=True, space_after=4)
        para(tf, b, size=10, color=MUTED, line=1.14)
    footnote(s, F["framework_foot"])
    notes(s, F["framework_notes"])

    # ======================================================= 12 GAPS
    s = slide(prs)
    y = heading(s, "What the framework makes visible",
                eyebrow="Findings · identified gaps",
                sub="Each gap is a thin or empty region of the matrix, paired with the "
                    "artifact hypothesised to close it.")
    _, tf = tb(s, M, y, CONTENT_W, Inches(0.30))
    para(tf, F["gaps_scope"], size=13, bold=True, color=CRIT, font=H_FONT, first=True)
    y += Inches(0.40)
    gaps = F["gaps"]
    gh = (H - Inches(0.86) - y - Inches(0.10) * (len(gaps) - 1)) / len(gaps)
    for i, g in enumerate(gaps):
        gy = y + i * (gh + Inches(0.10))
        core = g["severity"] == "core"
        card(s, M, gy, CONTENT_W, gh, fill=ACCENT_L if core else SOFT,
             line_col=ACCENT if core else LINE)
        badge(s, M + Inches(0.24), gy + Inches(0.14), Inches(0.40), g["id"],
              fill=ACCENT if core else MUTED, size=11)
        _, tf = tb(s, M + Inches(0.82), gy + Inches(0.10), Inches(3.30), gh - Inches(0.20),
                   anchor=MSO_ANCHOR.MIDDLE)
        para(tf, g["statement"], size=12.5, bold=True, color=NAVY, font=H_FONT, first=True,
             line=1.12)
        _, tf = tb(s, M + Inches(4.24), gy + Inches(0.10), Inches(4.55), gh - Inches(0.20),
                   anchor=MSO_ANCHOR.MIDDLE)
        para(tf, g["evidence_basis"], size=9.5, color=MUTED, first=True, line=1.14)
        _, tf = tb(s, M + Inches(8.95), gy + Inches(0.10), Inches(3.05), gh - Inches(0.20),
                   anchor=MSO_ANCHOR.MIDDLE)
        para(tf, g["which_sq"], size=10, bold=True, color=ACCENT if core else MUTED,
             first=True, line=1.14)
    footnote(s, "C1–C6 are design and research-contribution hypotheses. They may be called "
                "implemented only when the artifact exists and passes tests, and empirically "
                "beneficial only after independent evaluation.")
    notes(s, F["gaps_notes"])

    # ======================================================= 13 DIRECTION
    s = slide(prs, dark=True)
    y = heading(s, "The gap is integration — and here is what would falsify it",
                eyebrow="Conclusion · direction and hard gates", dark=True)
    card(s, M, y, CONTENT_W, Inches(1.50), fill=NAVY_SOFT)
    _, tf = tb(s, M + Inches(0.32), y + Inches(0.18), CONTENT_W - Inches(0.64), Inches(1.18))
    para(tf, F["chosen_gap"], size=15.5, bold=True, color=WHITE, font=H_FONT, first=True,
         line=1.16, space_after=5)
    para(tf, F["chosen_gap_why"], size=11, color=ACCENT_L, line=1.20)

    y2 = y + Inches(1.70)
    steps = [
        ("1", "Finalize RQ terminology", "supervisor decision log; D-RQ-01/02"),
        ("2", "Execute the locked searches", "QL-01–QL-05, currently 0 of 5"),
        ("3", "Collect independent expert labels", "EXP-005 at 0/24; ≥20 targeted"),
        ("4", "Evaluate SQ1 and SQ2 first", "burden, traceability, same-domain reuse"),
        ("5", "Then transfer, frozen-store", "matched control + blind target labels"),
        ("6", "Medical only after gates", "readiness 0/6; Plan B stays software"),
    ]
    sw = (CONTENT_W - Inches(0.22) * 2) / 3
    sh_ = Inches(1.06)
    for i, (n, t, b) in enumerate(steps):
        x = M + (i % 3) * (sw + Inches(0.22))
        cy = y2 + (i // 3) * (sh_ + Inches(0.18))
        card(s, x, cy, sw, sh_, fill=NAVY_SOFT)
        badge(s, x + Inches(0.20), cy + Inches(0.18), Inches(0.36), n, fill=ACCENT, size=11)
        _, tf = tb(s, x + Inches(0.66), cy + Inches(0.14), sw - Inches(0.86), Inches(0.80))
        para(tf, t, size=12.5, bold=True, color=WHITE, font=H_FONT, first=True,
             space_after=3, line=1.1)
        para(tf, b, size=10.5, color=RGBColor(0xBF, 0xCE, 0xDE), line=1.16)

    yb = y2 + 2 * (sh_ + Inches(0.18)) + Inches(0.06)
    _, tf = tb(s, M, yb, CONTENT_W, Inches(0.62))
    para(tf, "Human judgment should be requested selectively, captured with its reasoning, "
             "governed as a contestable lifecycle object, reused only as scoped advice, and "
             "evaluated across a frozen leakage boundary.",
         size=12.5, bold=True, color=ACCENT_L, first=True, line=1.22)
    para(tf, "Thank you — questions welcome.", size=13, bold=True, color=WHITE,
         space_before=8)
    notes(s, """
Pull it together. The gap I take is the connected one - not "when to ask", not "what to
store", but the fact that these are studied separately and the handovers between them are
where the open questions live.

Then the gates, in order, because a research agenda without gates is a wish list. Note that
three of the six are hard blockers I do not control: the wording decision, the labels, and
medical governance.

Close on the design thesis sentence. That is the one sentence I want them to remember.
(~55 sec)
""")

    # ======================================================= BACKUP 1
    s = slide(prs)
    y = heading(s, "Exact provisional wording and construct boundaries",
                eyebrow="Backup · terminology",
                sub="Preserved verbatim so the wording can be challenged precisely.")
    lw = (CONTENT_W - Inches(0.26)) / 2
    card(s, M, y, lw, Inches(3.94), fill=SOFT, line_col=LINE)
    _, tf = tb(s, M + Inches(0.26), y + Inches(0.20), lw - Inches(0.52), Inches(3.55))
    para(tf, "Provisional questions", size=14.5, bold=True, color=NAVY, font=H_FONT,
         first=True, space_after=8)
    rich(tf, [("U-RQ  ", {"bold": True, "size": 11.5, "color": ACCENT}),
              (U_RQ, {"size": 11, "color": MUTED})], line=1.18, space_after=7)
    for key, _t, body, _s in SQS:
        rich(tf, [(f"{key}  ", {"bold": True, "size": 11.5, "color": ACCENT}),
                  (body, {"size": 11, "color": MUTED})], line=1.18, space_after=7)

    rx = M + lw + Inches(0.26)
    card(s, rx, y, lw, Inches(3.94), fill=WHITE, line_col=LINE)
    _, tf = tb(s, rx + Inches(0.26), y + Inches(0.20), lw - Inches(0.52), Inches(3.55))
    para(tf, "Construct boundaries", size=14.5, bold=True, color=NAVY, font=H_FONT,
         first=True, space_after=8)
    for t, d in [
        ("Human judgment", "a context-bound assessment, correction, rationale or policy "
                           "decision by an authorized person — not a label."),
        ("Expert judgment", "authority tied to domain, language, pedagogical or governance "
                            "competence; partial and contestable."),
        ("Substantial variability", "a contextually justified alternative consistent with "
                                    "language semantics and domain logic."),
        ("Occasional variability", "an error, misconception or non-compliance — the term "
                                   "does not mean rare, which is a construct-validity problem "
                                   "flagged as G6."),
        ("Reusable judgment", "a governed record retrievable under explicit scope, authority, "
                              "version and validation conditions."),
    ]:
        rich(tf, [(f"{t}  ", {"bold": True, "size": 11.5, "color": INK}),
                  (d, {"size": 11, "color": MUTED})], line=1.18, space_after=7)
    notes(s, """
Use if asked for exact wording, or if someone challenges the substantial/occasional
terminology - which is a fair challenge, and it is why G6 exists.
""")

    # ======================================================= BACKUP 2
    s = slide(prs)
    y = heading(s, "Evidence levels and permitted claim language",
                eyebrow="Backup · claim discipline",
                sub="Every statement in the review is bounded by the maturity of the evidence "
                    "behind it.")
    rows = [
        ("FT-A", "Official or author-hosted full text reviewed for method, finding and "
                 "limitation.  Current count: 14.",
         "May support bounded detailed claims with a locator.", GOOD),
        ("FT-B", "Substantial authoritative excerpt or repository record reviewed; full "
                 "article not fully inspected.  Current count: 6.",
         "May support only the directly visible claim.", ACCENT),
        ("ID-S", "Bibliographic identity verified — title, venue, DOI or authoritative "
                 "record checked.",
         "Discovery and candidate relevance only.", MUTED),
        ("AB-S", "Title or abstract screened without sufficient full-text evidence.",
         "Search-space description only — never gap proof.", MUTED),
        ("Synthesis", "Reasoning that integrates sources or proposes a mechanism.",
         "Must be labelled synthesis, implication or research hypothesis.", CRIT),
    ]
    rh = (H - Inches(0.80) - y) / len(rows)
    for i, (lvl, defn, use, col) in enumerate(rows):
        ry = y + i * rh
        card(s, M, ry, CONTENT_W, rh - Inches(0.10), fill=SOFT, line_col=LINE)
        _, tf = tb(s, M + Inches(0.26), ry + Inches(0.06), Inches(1.35), rh - Inches(0.22),
                   anchor=MSO_ANCHOR.MIDDLE)
        para(tf, lvl, size=13.5, bold=True, color=col, font=H_FONT, first=True)
        _, tf = tb(s, M + Inches(1.75), ry + Inches(0.06), Inches(5.6), rh - Inches(0.22),
                   anchor=MSO_ANCHOR.MIDDLE)
        para(tf, defn, size=11, color=MUTED, first=True, line=1.16)
        _, tf = tb(s, M + Inches(7.55), ry + Inches(0.06), Inches(4.45), rh - Inches(0.22),
                   anchor=MSO_ANCHOR.MIDDLE)
        para(tf, use, size=11, bold=True, color=INK, first=True, line=1.16)
    notes(s, """
This is the slide that answers "how do you know you are not overclaiming". Each claim in the
review is tagged to the maturity of its evidence, and the permitted wording follows from
that tag rather than from how confident I feel.
""")

    # ======================================================= BACKUP 3
    s = slide(prs)
    y = heading(s, "SQ1 — selective intervention is a resource-allocation problem",
                eyebrow="Backup · provisional answer",
                sub="Uncertainty alone is insufficient. Priority must also weigh consequence, "
                    "novelty, disagreement, reuse value and burden.")
    lw = (CONTENT_W - Inches(0.26)) / 2
    card(s, M, y, lw, Inches(3.80), fill=SOFT, line_col=LINE)
    _, tf = tb(s, M + Inches(0.28), y + Inches(0.22), lw - Inches(0.56), Inches(3.4))
    para(tf, "Trigger signals", size=14.5, bold=True, color=NAVY, font=H_FONT, first=True,
         space_after=9)
    for t in ["Calibrated uncertainty", "Cross-agent disagreement",
              "Novelty or coverage gap", "Consequence and policy importance",
              "Weak underlying evidence", "Expected value of information",
              "Reviewer burden and queue state"]:
        para(tf, "•  " + t, size=12, color=MUTED, line=1.2, space_after=6)

    rx = M + lw + Inches(0.26)
    card(s, rx, y, lw, Inches(3.80), fill=WHITE, line_col=LINE)
    _, tf = tb(s, rx + Inches(0.28), y + Inches(0.22), lw - Inches(0.56), Inches(3.4))
    para(tf, "Intervention modes", size=14.5, bold=True, color=NAVY, font=H_FONT, first=True,
         space_after=9)
    for t, d in [("Proceed and log", "low impact, familiar"),
                 ("Focused review", "uncertain but bounded"),
                 ("Queue or batch review", "important, not urgent"),
                 ("Interrupt or stop", "high consequence")]:
        rich(tf, [(f"{t}  ", {"bold": True, "size": 12, "color": INK}),
                  (d, {"size": 12, "color": MUTED})], line=1.2, space_after=8)
    para(tf, "Provisional answer: request judgment when the expected value of expert "
             "information exceeds review burden and interruption risk — targeting the "
             "smallest decision unit.",
         size=11.5, color=NAVY, italic=True, line=1.22, space_before=6)
    notes(s, """
Use if asked why a confidence threshold is not enough. The point is bounded expert attention:
a policy, not a threshold - and the priority score is a design hypothesis, not an estimated
model.
""")

    # ======================================================= BACKUP 4
    s = slide(prs)
    y = heading(s, "SQ2 — a reusable judgment is not a naked label",
                eyebrow="Backup · the governed judgment object",
                sub="The minimum reusable unit carries its grounding, its reasoning, its "
                    "authority and its lifecycle.")
    fields = [
        ("Case grounding", "artifact · fragment · guideline version · evidence locator"),
        ("Decision trace", "system claim · confidence · concise reasoning · correction"),
        ("Governance", "authority · scope · visibility · dissent · permission"),
        ("Lifecycle", "validated · expired · superseded · revoked · active"),
        ("Reuse signals", "semantic signature · counterexamples · thresholds"),
        ("Outcome trace", "retrieval · influence · override · downstream effect"),
    ]
    cw6 = (CONTENT_W - Inches(0.22) * 2) / 3
    for i, (t, d) in enumerate(fields):
        x = M + (i % 3) * (cw6 + Inches(0.22))
        cy = y + (i // 3) * (Inches(1.30) + Inches(0.20))
        card(s, x, cy, cw6, Inches(1.30), fill=SOFT, line_col=LINE)
        badge(s, x + Inches(0.22), cy + Inches(0.20), Inches(0.38), str(i + 1), size=11)
        _, tf = tb(s, x + Inches(0.22), cy + Inches(0.70), cw6 - Inches(0.44), Inches(0.5))
        para(tf, t, size=13, bold=True, color=NAVY, font=H_FONT, first=True, space_after=3)
        para(tf, d, size=10.5, color=MUTED, line=1.16)
    yb = y + 2 * Inches(1.50) + Inches(0.10)
    card(s, M, yb, CONTENT_W, Inches(0.92), fill=NAVY)
    _, tf = tb(s, M + Inches(0.30), yb + Inches(0.16), CONTENT_W - Inches(0.60), Inches(0.62))
    para(tf, "Reasoning-capture rule", size=12.5, bold=True, color=ACCENT_L, first=True,
         space_after=4)
    para(tf, "Store the shortest sufficient decision trace — claim, evidence, rule, "
             "uncertainty and correction. Not hidden chain-of-thought, and not unnecessary "
             "personal data.",
         size=12, color=WHITE, line=1.2)
    notes(s, """
Use if asked what "capturing judgment" actually means. The acceptance test is
reconstruction: could an authorized reviewer who was not present rebuild why this ruling was
made? The reasoning-capture rule at the bottom is the privacy and proportionality boundary.
""")

    # ======================================================= BACKUP 5
    s = slide(prs)
    y = heading(s, "SQ3 — retrieval is not permission",
                eyebrow="Backup · transfer ladder and leakage controls",
                sub="Reuse is advisory and scope-filtered; evidence requirements rise with "
                    "transfer distance.")
    ladder = [
        ("L0", "same-case retry", "default relevance high", GOOD),
        ("L1", "same task and domain", "advisory after scope checks", GOOD),
        ("L2", "new task, same domain", "independent labels + holdout", ORANGE),
        ("L3", "adjacent domain", "preregistered restricted experiment", ORANGE),
        ("L4", "cross-domain / healthcare", "no default permission; ethics + authorization", CRIT),
    ]
    lh = Inches(0.66)
    for i, (lvl, what, need, col) in enumerate(ladder):
        cy = y + i * (lh + Inches(0.10))
        card(s, M, cy, Inches(7.9), lh, fill=SOFT, line_col=LINE)
        badge(s, M + Inches(0.18), cy + Inches(0.11), Inches(0.44), lvl, fill=col, size=11)
        _, tf = tb(s, M + Inches(0.82), cy + Inches(0.06), Inches(2.9), lh - Inches(0.12),
                   anchor=MSO_ANCHOR.MIDDLE)
        para(tf, what, size=12.5, bold=True, color=NAVY, font=H_FONT, first=True)
        _, tf = tb(s, M + Inches(3.85), cy + Inches(0.06), Inches(3.9), lh - Inches(0.12),
                   anchor=MSO_ANCHOR.MIDDLE)
        para(tf, need, size=11.5, color=MUTED, first=True, line=1.14)

    rx = M + Inches(8.12)
    rw = CONTENT_W - Inches(8.12)
    card(s, rx, y, rw, Inches(3.70), fill=RGBColor(0xFD, 0xF0, 0xF0),
         line_col=RGBColor(0xF0, 0xC8, 0xC8))
    _, tf = tb(s, rx + Inches(0.26), y + Inches(0.20), rw - Inches(0.52), Inches(3.3))
    para(tf, "Leakage-safe evaluation", size=14, bold=True, color=CRIT, font=H_FONT,
         first=True, space_after=9)
    for t in ["Frozen source store and retrieval policy",
              "Matched no-reuse control",
              "Independent blind target labels, set before source judgments",
              "Store updated only after scoring",
              "Predeclared thresholds"]:
        para(tf, "•  " + t, size=11.5, color=MUTED, line=1.2, space_after=7)
    notes(s, """
Use for the transfer question. The one line to land: similarity retrieves, permission
governs. And updating the store from target labels before evaluation would break the
frozen-store design and void any transfer evidence - which is why it matters.
""")

    # ======================================================= BACKUP 6
    s = slide(prs)
    y = heading(s, "Three studies, each independently falsifiable",
                eyebrow="Backup · research contract",
                sub="Separate units of analysis prevent one small label set from carrying "
                    "every claim.")
    studies = [
        ("Study 1 / SQ1", "Selective Human Review Orchestrator",
         "Held-out cases, policy baselines, burden logging",
         "Falsified by: no gain over a simpler baseline, or unacceptable burden"),
        ("Study 2 / SQ2", "Governed Judgment Object + Contestable Store",
         "Label-only comparator, reconstruction and lifecycle tests",
         "Falsified by: reviewers cannot use the object, or governance cost outweighs benefit"),
        ("Study 3 / SQ3", "Scope-Aware Retrieval Advisor + transfer classifier",
         "Frozen store, blind labels, matched no-reuse control",
         "Falsified by: no target benefit, or effect vanishes under blinding"),
        ("Integrated / U-RQ", "End-to-end governed lifecycle",
         "Human-only, AI-only and non-governed baselines",
         "Falsified by: any load-bearing study fails, or gain is only extra human effort"),
    ]
    sh2 = (H - Inches(0.80) - y) / len(studies)
    for i, (name, artifact, design, falsify) in enumerate(studies):
        cy = y + i * sh2
        card(s, M, cy, CONTENT_W, sh2 - Inches(0.12), fill=SOFT, line_col=LINE)
        _, tf = tb(s, M + Inches(0.24), cy + Inches(0.08), Inches(2.05), sh2 - Inches(0.28),
                   anchor=MSO_ANCHOR.MIDDLE)
        para(tf, name, size=12, bold=True, color=ACCENT, font=H_FONT, first=True, line=1.12)
        _, tf = tb(s, M + Inches(2.42), cy + Inches(0.08), Inches(3.5), sh2 - Inches(0.28),
                   anchor=MSO_ANCHOR.MIDDLE)
        para(tf, artifact, size=11, bold=True, color=NAVY, first=True, line=1.14)
        _, tf = tb(s, M + Inches(6.05), cy + Inches(0.08), Inches(3.0), sh2 - Inches(0.28),
                   anchor=MSO_ANCHOR.MIDDLE)
        para(tf, design, size=10.5, color=MUTED, first=True, line=1.14)
        _, tf = tb(s, M + Inches(9.20), cy + Inches(0.08), Inches(2.85), sh2 - Inches(0.28),
                   anchor=MSO_ANCHOR.MIDDLE)
        para(tf, falsify, size=10, color=CRIT, first=True, line=1.14)
    notes(s, """
Use if asked "how would you know you are wrong". Every study has a stated falsification
condition, written before the data exists. Note also that EXP-005 alone cannot satisfy all
three - each has its own estimand.
""")

    # ======================================================= BACKUP 7
    s = slide(prs)
    y = heading(s, "Design science, GenAI use, and the evidence boundary",
                eyebrow="Backup · method and honesty controls",
                sub="How this sits in the course's design-science frame, and exactly what is "
                    "not being claimed.")
    cw7 = (CONTENT_W - Inches(0.22) * 2) / 3
    card(s, M, y, cw7, Inches(3.86), fill=SOFT, line_col=LINE)
    _, tf = tb(s, M + Inches(0.24), y + Inches(0.20), cw7 - Inches(0.48), Inches(3.5))
    para(tf, "Design science", size=14, bold=True, color=ACCENT, font=H_FONT, first=True,
         space_after=9)
    for t, d in [("Relevance", "expert review is scarce; deviation meaning is contextual; "
                               "guidelines drift unseen"),
                 ("Design", "selective-intervention policy + governed judgment lifecycle, "
                            "built and evaluated per sub-question"),
                 ("Rigor", "structured review with evidence levels; the lifecycle framework "
                           "is itself a candidate contribution back")]:
        rich(tf, [(f"{t}  ", {"bold": True, "size": 11.5, "color": INK}),
                  (d, {"size": 11, "color": MUTED})], line=1.18, space_after=8)

    x2 = M + cw7 + Inches(0.22)
    card(s, x2, y, cw7, Inches(3.86), fill=RGBColor(0xF3, 0xF8, 0xF3),
         line_col=RGBColor(0xCC, 0xE4, 0xCC))
    _, tf = tb(s, x2 + Inches(0.24), y + Inches(0.20), cw7 - Inches(0.48), Inches(3.5))
    para(tf, "GenAI — used for", size=14, bold=True, color=GOOD, font=H_FONT, first=True,
         space_after=8)
    for t in ["Literature organisation", "Synthesis drafting", "Diagram production",
              "Consistency checks", "Document layout"]:
        para(tf, "•  " + t, size=11.5, color=MUTED, line=1.18, space_after=5)
    para(tf, "Not used for", size=13, bold=True, color=CRIT, space_before=8, space_after=6)
    for t in ["Approving inclusion decisions", "Inventing citations",
              "Creating expert labels", "Claiming supervisor approval"]:
        para(tf, "•  " + t, size=11.5, color=MUTED, line=1.18, space_after=5)

    x3 = x2 + cw7 + Inches(0.22)
    card(s, x3, y, cw7, Inches(3.86), fill=RGBColor(0xFD, 0xF0, 0xF0),
         line_col=RGBColor(0xF0, 0xC8, 0xC8))
    _, tf = tb(s, x3 + Inches(0.24), y + Inches(0.20), cw7 - Inches(0.48), Inches(3.5))
    para(tf, "Not claimed", size=14, bold=True, color=CRIT, font=H_FONT, first=True,
         space_after=9)
    for t in ["No improvement in assessment accuracy — expert labels stand at 0 of 24.",
              "No reduction in expert workload.",
              "No safe generalization or transfer.",
              "Nothing about clinical performance — medical readiness 0 of 6.",
              "No claim that the literature holds no prior solution — the formal "
              "searches are not yet executed."]:
        para(tf, "•  " + t, size=11, color=MUTED, line=1.18, space_after=7)
    notes(s, """
The slide I would use if someone asks "so does it work?" or "how much of this did AI write?".

Honest answer to the first: I can show mechanism and observability and the trade-off. I
cannot show quality, effort or transfer, because that needs independent expert labels and I
have none validated.

Honest answer to the second: AI helped organise, draft and draw. It did not decide what to
include, invent a citation, or create a label. Stating that boundary is what makes the rest
credible.
""")

    prs.save(DECK)
    n = len(prs.slides._sldIdLst)
    print(f"saved: {DECK}")
    print(f"slides: {n}  (13 presented + 7 backup)")
    assert n <= 20, f"deck exceeds the 20-slide limit: {n}"
    return DECK


if __name__ == "__main__":
    build()
