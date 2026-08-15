#!/usr/bin/env python3
"""Build the Aug-12 progress presentation (deliverable A08-08).

Deliberately plain and honest: no claim appears on a slide that is not supported
by the package, and the evidence-gate slide states the blocking counts directly.
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ACCENT = RGBColor(0x1F, 0x38, 0x64)
INK = RGBColor(0x25, 0x25, 0x25)
MUTED = RGBColor(0x60, 0x60, 0x60)
WARN = RGBColor(0xB0, 0x4A, 0x00)
GOOD = RGBColor(0x2E, 0x6B, 0x2E)
RULE = RGBColor(0xD0, 0xD7, 0xE2)


def add_slide(prs: Presentation, title: str, kicker: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if kicker:
        box = slide.shapes.add_textbox(Inches(0.6), Inches(0.32), Inches(12.1), Inches(0.3))
        p = box.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = kicker.upper()
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = MUTED
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.62), Inches(12.1), Inches(0.75))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = ACCENT
    line = slide.shapes.add_shape(1, Inches(0.6), Inches(1.42), Inches(12.1), Emu(12700))
    line.fill.solid()
    line.fill.fore_color.rgb = RULE
    line.line.fill.background()
    line.shadow.inherit = False
    return slide


def bullets(slide, items, top=1.75, size=17, left=0.75, width=11.9):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.0))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        text, level, color = (item + (None,))[:3] if isinstance(item, tuple) else (item, 0, None)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(9)
        r = p.add_run()
        r.text = ("• " if level == 0 else "– ") + text
        r.font.size = Pt(size if level == 0 else size - 2)
        r.font.color.rgb = color or (INK if level == 0 else MUTED)
    return box


def table(slide, headers, rows, top=1.8, left=0.7, width=12.0, height=1.0, fontsize=12):
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tbl = shape.table
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(fontsize)
                r.font.bold = True
                r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
    for i, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = tbl.cell(i, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(fontsize)
                    r.font.color.rgb = INK
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if i % 2 else RGBColor(0xF2, 0xF5, 0xFA)
    return tbl


def footer(slide, text):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.4))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(10.5)
    r.font.italic = True
    r.font.color.rgb = MUTED


def build(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1 title
    s = prs.slides.add_slide(prs.slide_layouts[6])
    box = s.shapes.add_textbox(Inches(0.9), Inches(2.3), Inches(11.5), Inches(1.4))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Chapter 3 Complete — Gap and Research Questions"
    r.font.size = Pt(40)
    r.font.bold = True
    r.font.color.rgb = ACCENT
    box2 = s.shapes.add_textbox(Inches(0.95), Inches(3.7), Inches(11.5), Inches(1.6))
    tf = box2.text_frame
    for i, line in enumerate([
        "PhD-track proposal · progress review",
        "Ali Hamed  ·  Prof. Iris Reinhartz-Berger  ·  Prof. Arnon",
        "Wednesday 12 August 2026, 09:00",
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = line
        r.font.size = Pt(17 if i == 0 else 15)
        r.font.color.rgb = INK if i else MUTED
    footer(s, "Working drafts for review — not supervisor-approved. No accuracy, generalization or clinical claim is made.")

    # 2 what you asked for
    s = add_slide(prs, "What you asked for on 5 August", "agenda item 1")
    table(s, ["Request (E15)", "State"], [
        ["Finish the Gap and Research Questions chapter in full", "Delivered"],
        ["Per-RQ literature spreadsheet with an RQ tag column", "Delivered"],
        ["Think about §2 and §4 — do not start them", "Thought about, not started"],
        ["Word proposal + a separate tracking document", "Delivered (both)"],
        ["Share the Drive with Iris and Arnon", "Needs Ali (Google account)"],
        ["Reply to your check-in email", "Awaiting the email"],
        ["Present progress live today", "This session"],
    ], top=1.85, height=3.9, fontsize=14)
    footer(s, "Five of seven delivered in the folder; two require Ali personally.")

    # 3 corrections applied
    s = add_slide(prs, "Your corrections — all applied, two flagged back", "agenda item 2")
    bullets(s, [
        ("E5 — \"reused\" removed from the main question; reuse now lives in SQ2 and SQ3", 0, GOOD),
        ("E6 — question narrowed to variability in guideline operationalization", 0, GOOD),
        ("E7 — \"reliable\" kept; auditable / transferable / end-to-end dropped from the headline", 0, GOOD),
        ("E8 — SQ2 split into capture vs. transfer; \"expert judgment\" adopted", 0, GOOD),
        ("E9 — SQ2 tied to core reasoning; evaluation criteria built into all three questions", 0, GOOD),
        ("E12 — SQ3 rebuilt around classifying domain-specific vs. transferable uncertainty", 0, GOOD),
        ("E13 — proposal follows the confirmed chapter structure", 0, GOOD),
        ("Two things I could NOT resolve from the recording — I did not choose silently:", 0, WARN),
        ("\"exploration\" vs \"identification / classification\" in the main question", 1, WARN),
        ("\"human judgment\" vs \"expert judgment\" across the question set", 1, WARN),
    ], top=1.8, size=15)
    footer(s, "Both need my saved working draft (A08-01) and your ruling.")

    # 4 the gap
    s = add_slide(prs, "The gap: four fields, each solving a different part", "agenda item 3 — Chapter 3 §3.2")
    table(s, ["Body of work", "Establishes", "Leaves open"], [
        ["Automated model assessment", "Deviation can be detected and scored", "What a deviation MEANS"],
        ["LLM / agentic assistance", "Expert correction is consistently needed", "WHEN to ask; WHAT to keep"],
        ["Human-in-the-loop & oversight", "Mechanisms and governance expectations", "Selection policy under a budget"],
        ["Variability engineering", "Formalism for DESIGNED variability", "OBSERVED interpretive variability"],
    ], top=1.9, height=3.0, fontsize=14)
    b = s.shapes.add_textbox(Inches(0.75), Inches(5.15), Inches(11.9), Inches(1.4))
    p = b.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = ("The gap is the CONNECTION between them:  when to ask  →  what to keep  →  "
              "where it may legitimately be reused.")
    r.font.size = Pt(17)
    r.font.bold = True
    r.font.color.rgb = ACCENT
    footer(s, "Openness statements are candidate claims — the systematic searches are frozen and not yet run.")

    # 5 questions
    s = add_slide(prs, "One question, three sub-questions", "agenda item 4 — Chapter 3 §3.4–§3.7")
    b = s.shapes.add_textbox(Inches(0.75), Inches(1.75), Inches(11.9), Inches(1.1))
    tf = b.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = ("U-RQ  ·  How can human judgment be captured, governed, and used to support agentic-AI-driven "
              "variability exploration in guideline operationalization scenarios, enabling reliable human–AI "
              "co-reasoning?")
    r.font.size = Pt(15)
    r.font.italic = True
    r.font.color.rgb = ACCENT
    table(s, ["", "Asks", "Built-in evaluation criterion"], [
        ["SQ1", "WHEN should the system ask an expert?", "Coverage of important uncertainty AND bounded burden"],
        ["SQ2", "WHAT must be kept, under what governance?", "No unsafe generalization; no loss of human authority"],
        ["SQ3", "WHERE ELSE does a judgment hold?", "Classification must discriminate; transfer must survive leakage controls"],
    ], top=3.05, height=2.3, fontsize=13)
    footer(s, "Each question can fail on its own terms — and a negative result is still a reportable contribution.")

    # 6 literature
    s = add_slide(prs, "Literature: 40 sources tagged — and what they do NOT cover", "agenda item 5")
    bullets(s, [
        ("40 verified sources, each tagged RQ1 / RQ2 / RQ3 / general, independently fabrication-checked", 0),
        ("Coverage-gap analysis per question — the part beyond an inventory:", 0),
        ("RQ1 — nothing on interrupt/approval mechanics inside a multi-agent LLM pipeline", 1, WARN),
        ("RQ2 — weakest coverage: no agent-memory or learning-from-feedback sources", 1, WARN),
        ("RQ3 — no usability instruments, no cross-domain transfer studies", 1, WARN),
        ("general — NO design-science methodology sources, though the method is design science", 1, WARN),
        ("Frozen search protocol QL-01…QL-05 ready to execute on your word", 0),
    ], top=1.85, size=16)
    footer(s, "Seed corpus from our own MAS4Models reference list + the tracked resource pack. Not a systematic review.")

    # 7 evidence
    s = add_slide(prs, "Evidence position — unchanged, and stated plainly", "agenda item 5")
    table(s, ["Gate", "Value", "What it blocks"], [
        ["Independent expert labels (EXP-005)", "0 of 24", "Every quantitative accuracy claim (≥20 needed)"],
        ["Medical entry gates G1–G6", "0 of 6", "Any medical data processing"],
        ["Literature searches QL-01…QL-05", "Not run", "Any novelty or completeness statement"],
        ["Accuracy / generalization / clinical claims", "None made", "— checked by an automated guard"],
    ], top=1.95, height=2.6, fontsize=14)
    b = s.shapes.add_textbox(Inches(0.75), Inches(4.85), Inches(11.9), Inches(1.5))
    tf = b.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = ("The binding constraint is not writing or engineering — it is independent expert labels. "
              "That needs two reviewers and an adjudicator: a resourcing decision, not a build task.")
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = WARN
    footer(s, "Nothing in this package depends on evidence we do not have.")

    # 8 sections 2 and 4
    s = add_slide(prs, "§2 and §4 — thought about, deliberately not started", "agenda item 6")
    bullets(s, [
        ("§2 literature survey: four structural options with trade-offs (SQ-parallel / community-parallel / "
         "lifecycle-stage / hybrid) — none selected", 0),
        ("§4 research artifact: nine options across SQ1–SQ3, at different abstraction levels — none selected", 0),
        ("A recurring issue surfaced: the current \"core artifact\" rows bundle 6, 9 and 10 components; "
         "several SQ3 items are study OUTPUTS, not artifacts", 0, WARN),
        ("14 open questions prepared for you — grouped by chapter shape, artifact definition, "
         "evidence admissibility, and resourcing", 0),
    ], top=1.9, size=16)
    footer(s, "Compliance with your instruction — the thinking is delivered; neither section is begun.")

    # 9 decisions
    s = add_slide(prs, "Eight decisions I need from you", "agenda item 7")
    table(s, ["ID", "Decision", "Blocks"], [
        ["D-RQ-01", "Final U-RQ and SQ1–SQ3 wording (incl. the two flagged discrepancies)", "Everything downstream"],
        ["D-RQ-02", "Ratify or reverse dropping auditable / transferable / end-to-end", "Title and headline"],
        ["D-TITLE-01", "Proposal title — three candidates offered", "Submission-readiness"],
        ["D-DECOMP-01", "Are three sub-questions jointly sufficient for \"reliable co-reasoning\"?", "Chapter 3 sign-off"],
        ["D-CH2-01", "Chapter 2 structure — which of four options", "Starting §2"],
        ["D-ART-01", "One named artifact per study, or a package?", "Starting §4"],
        ["D-SEARCH-01", "Run searches before or after wording sign-off? Add QL-06+?", "Literature execution"],
        ["D-PRELIM-01", "Which preliminary results may appear, labelled how?", "Chapter 5"],
    ], top=1.9, height=4.0, fontsize=12)
    footer(s, "The first three block everything else.")

    # 10 close
    s = add_slide(prs, "Where this leaves us")
    bullets(s, [
        ("Chapter 3 is complete, cited, and honest about what is still a candidate claim", 0, GOOD),
        ("The literature map exists — and its gaps are now visible rather than latent", 0, GOOD),
        ("§2 and §4 are ready to start the moment you confirm direction and wording", 0),
        ("Two wording questions and eight decisions are on the table today", 0, WARN),
        ("Next after this meeting: apply your rulings in one pass, then execute the searches", 0),
    ], top=1.95, size=17)
    footer(s, "All materials in the shared Drive folder: 2026-08-12 Supervisor Package.")

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"presentation: {len(prs.slides.__iter__.__self__._sldIdLst)} slides -> {out}")


if __name__ == "__main__":
    build(Path(sys.argv[1]) if len(sys.argv) > 1 else Path("outputs/aug12/04 - Progress Presentation.pptx"))
